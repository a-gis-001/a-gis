def _read_to_text_pptx(*, file: type["pathlib.Path"]) -> str:
    """Extracts text from a PPTX file.

    This function reads a PPTX file and extracts all text from its slides,
    returning it as a single string with each piece of text separated by a newline.

    Args:
        file (pathlib.Path): The path to the PPTX file to be read.

    Returns:
        str: The extracted text from the PPTX file.
    """
    import pptx
    import re
    import pathlib

    # Load the presentation
    presentation = pptx.Presentation(file)
    extracted_text = []

    # Iterate through each slide
    for slide_num, slide in enumerate(presentation.slides, start=1):
        slide_text = [f"Slide {slide_num}:\n"]  # Add slide number header

        # Iterate through all shapes within the slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Extract paragraphs from the shape's text frame
                for paragraph in shape.text_frame.paragraphs:
                    # Combine all runs in the paragraph
                    para_text = " ".join(run.text for run in paragraph.runs)
                    slide_text.append(para_text)

        # Append the slide's text, joined by newlines
        extracted_text.append("\n".join(slide_text))

    # Join all slides' text with double newlines between slides
    text = "\n\n".join(extracted_text)

    # Use regex to clean unwanted characters, if needed (e.g., removing '^K')
    # text = re.sub(r"\^K", " ", text_output)

    return text
